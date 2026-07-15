import json
import tempfile
import unittest
from pathlib import Path

from canon.config import load_settings
from canon.ingest.flexible import ingest_flexible_source, profile_source


class FlexibleIngestTests(unittest.TestCase):
    def test_profile_detects_crm_csv_shape(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "crm.csv"
            path.write_text(
                "id,account_id,opportunity_stage,notes,owner\n"
                "n1,acme,evaluation,Customer needs better evidence review,Ada\n",
                encoding="utf-8",
            )

            profile = profile_source(path)

        self.assertEqual(profile["detected_format"], "csv")
        self.assertEqual(profile["source_shape"], "crm_record")
        self.assertEqual(profile["proposed_mapping"]["document_type"], "crm_record")
        self.assertEqual(profile["chunking_strategy"], "field_aware_record")

    def test_ingest_csv_normalizes_records_and_writes_processed_files(self):
        settings = load_settings()
        mode = "unit_flexible_csv"
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "crm.csv"
            path.write_text(
                "id,account_id,opportunity_stage,notes,owner\n"
                "n1,acme,evaluation,Customer asked for renewal-risk evidence,Ada\n"
                "n2,bravo,proposal,,Ben\n",
                encoding="utf-8",
            )

            report = ingest_flexible_source(path, mode=mode, chunk_tokens=24, overlap_tokens=4)

        self.assertEqual(report["raw_record_count"], 2)
        self.assertEqual(report["normalized_record_count"], 1)
        self.assertEqual(report["skipped_reasons"]["missing_text"], 1)
        self.assertGreaterEqual(report["chunk_count"], 1)
        works_path = settings.data_dir / "processed" / f"works_{mode}.json"
        works = json.loads(works_path.read_text(encoding="utf-8"))
        self.assertEqual(works[0]["raw"]["document_type"], "crm_record")
        self.assertEqual(works[0]["raw"]["source_shape"], "crm_record")

    def test_ingest_text_folder_profiles_mixed_source(self):
        settings = load_settings()
        mode = "unit_flexible_folder"
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            (root / "note.md").write_text("# Topic\n\nInterconnection delays affect delivery.", encoding="utf-8")
            (root / "memo.txt").write_text("Financing costs may delay storage projects.", encoding="utf-8")

            report = ingest_flexible_source(root, mode=mode, chunk_tokens=24, overlap_tokens=4)

        self.assertEqual(report["profile"]["source_shape"], "mixed_source")
        self.assertEqual(report["normalized_record_count"], 2)
        self.assertGreaterEqual(report["chunk_count"], 2)
        raw_path = settings.data_dir / "raw" / f"unstructured_{mode}.json"
        self.assertTrue(raw_path.exists())

    def test_profile_research_json_records(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "papers.json"
            path.write_text(
                json.dumps(
                    {
                        "records": [
                            {
                                "id": "paper-1",
                                "title": "Evidence Systems",
                                "authors": ["A. Scholar"],
                                "doi": "10.123/example",
                                "abstract": "We study evidence retrieval systems.",
                            }
                        ]
                    }
                ),
                encoding="utf-8",
            )

            profile = profile_source(path)

        self.assertEqual(profile["source_shape"], "research_paper")
        self.assertEqual(profile["proposed_mapping"]["document_type"], "academic_article")
        self.assertEqual(profile["proposed_mapping"]["confidence"], "high")

    def test_ingest_docx_html_and_xlsx_files(self):
        try:
            from docx import Document
            from openpyxl import Workbook
        except ImportError as exc:
            self.skipTest(f"Optional document parser missing: {exc}")
        settings = load_settings()
        mode = "unit_flexible_rich_files"
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            document = Document()
            document.add_paragraph("AI data centers can affect local water demand.")
            document.save(root / "brief.docx")
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "projects"
            sheet.append(["region", "risk"])
            sheet.append(["Chile", "water stress"])
            workbook.save(root / "projects.xlsx")
            (root / "page.html").write_text(
                "<html><head><title>Grid Memo</title></head><body><p>Grid interconnection delays affect data center buildout.</p></body></html>",
                encoding="utf-8",
            )

            report = ingest_flexible_source(root, mode=mode, chunk_tokens=24, overlap_tokens=4)

        self.assertEqual(report["normalized_record_count"], 3)
        self.assertGreaterEqual(report["chunk_count"], 3)
        works_path = settings.data_dir / "processed" / f"works_{mode}.json"
        works = json.loads(works_path.read_text(encoding="utf-8"))
        document_types = {work["raw"]["document_type"] for work in works}
        self.assertIn("word_document", document_types)
        self.assertIn("spreadsheet", document_types)
        self.assertIn("html_document", document_types)

    def test_ingest_code_and_notebook_files(self):
        settings = load_settings()
        mode = "unit_flexible_code_files"
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            (root / "pipeline.py").write_text(
                "def build_corpus(path):\n    return f'index {path}'\n",
                encoding="utf-8",
            )
            (root / "config.yaml").write_text(
                "retrieval:\n  top_k: 10\n  method: hybrid\n",
                encoding="utf-8",
            )
            (root / "analysis.ipynb").write_text(
                json.dumps(
                    {
                        "cells": [
                            {"cell_type": "markdown", "source": ["# Evidence notes\n"]},
                            {"cell_type": "code", "source": ["query = 'data center water risk'\n"]},
                        ]
                    }
                ),
                encoding="utf-8",
            )

            report = ingest_flexible_source(root, mode=mode, chunk_tokens=24, overlap_tokens=4)

        self.assertEqual(report["normalized_record_count"], 3)
        works_path = settings.data_dir / "processed" / f"works_{mode}.json"
        works = json.loads(works_path.read_text(encoding="utf-8"))
        document_types = {work["raw"]["document_type"] for work in works}
        self.assertIn("source_code", document_types)
        self.assertIn("notebook", document_types)

    def test_ingest_pptx_presentation_files(self):
        try:
            from pptx import Presentation
        except ImportError as exc:
            self.skipTest(f"Optional presentation parser missing: {exc}")
        settings = load_settings()
        mode = "unit_flexible_presentation_files"
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[5])
            slide.shapes.title.text = "AI infrastructure risks"
            text_box = slide.shapes.add_textbox(0, 0, 4000000, 1000000)
            text_box.text = "Power availability and water stress affect deployment."
            deck.save(root / "briefing.pptx")

            report = ingest_flexible_source(root, mode=mode, chunk_tokens=24, overlap_tokens=4)

        self.assertEqual(report["normalized_record_count"], 1)
        works_path = settings.data_dir / "processed" / f"works_{mode}.json"
        works = json.loads(works_path.read_text(encoding="utf-8"))
        self.assertEqual(works[0]["raw"]["document_type"], "presentation")

    def test_google_pointer_and_image_are_detected_but_not_ingested_as_text(self):
        try:
            from PIL import Image
        except ImportError as exc:
            self.skipTest(f"Optional image parser missing: {exc}")
        mode = "unit_flexible_non_text_files"
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            (root / "doc.gdoc").write_text(
                json.dumps({"name": "Native Doc", "url": "https://docs.google.com/document/d/example"}),
                encoding="utf-8",
            )
            Image.new("RGB", (10, 10), color="white").save(root / "scan.png")

            profile = profile_source(root)
            report = ingest_flexible_source(root, mode=mode, chunk_tokens=24, overlap_tokens=4)

        self.assertIn("Some files were detected", " ".join(profile["warnings"]))
        self.assertEqual(report["raw_record_count"], 2)
        self.assertEqual(report["normalized_record_count"], 0)
        self.assertEqual(report["skipped_reasons"]["missing_text"], 2)


if __name__ == "__main__":
    unittest.main()
